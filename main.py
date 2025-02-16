import json
import requests
import sys
from dotenv import load_dotenv
import os
from io import BytesIO
from typing import Dict, List
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
import subprocess
from pydub import AudioSegment
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips, VideoFileClip
import logging
from tqdm import tqdm

from image_gen.image_gen import generate_image

load_dotenv()

# # Configure logging
# logging.basicConfig(
#     level=logging.INFO,
#     format='%(asctime)s - %(levelname)s - %(message)s'
# )


# load word list from file
with open("materials/common_3000_eng_words.txt", "r") as file:
    word_list = file.read().splitlines()

output_dir = os.getcwd() + "/output/"
cache_dir = os.getcwd() + "/cache/"
# Anthropic
import anthropic
anthropic_client = anthropic.Anthropic()

# OpenAI
import openai
openai_client = openai.OpenAI()


def generate_word_info(word: str) -> Dict[str, str]:
    """
    Generate a word's definition and example sentence from a given word.

    Args:
        word (str): The word to generate information for.

    Returns:
        dict: A dictionary containing the word's definition and example sentence under the keys 'definition' and 'example'.
    """
    system_message: str = "You are a kintergarden teacher. You are given a word and you need to explain it in a way that is easy to understand."
    user_message: str = (f"Can you explain the word {word} in a way that is easy to understand for a 5 year old? "
                f"Please respond in JSON format with the following two fields: 'definition' and 'example'. "
                f"The definition should be no more than a couple of sentences and no more than 25 words explaining the most common definition(s) of the word. "
                f"The example should be a sentence and no more than 25 words that uses the word in a way that is easy to understand for a 5 year old. "
                f"Start the example sentence with something like 'For example, ...', 'Here is an example: ...', or 'An example would be ...'."
                f"\n\n"
                f"Here is an example JSON response: "
                f"{{"
                f"  'definition': 'A definition of the word', "
                f"  'example': 'An example sentence using the word'"
                f"}}")
    
    message: anthropic.Message = anthropic_client.messages.create(
        model="claude-3-5-sonnet-20241022",
        max_tokens=1024,
        temperature=0,
        system=system_message,
        messages=[
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": user_message
                }
            ]
            }
        ]
    )
    # print(message.content)

    try:
        word_info: Dict[str, str] = json.loads(message.content[0].text)
        word_info["word"] = word
        return word_info
    except:
        print("Failed to parse JSON response for word: " + word + " message: " + message.content[0].text)
        # record word and message to file
        with open(output_dir + "failed_words.txt", "a") as file:
            file.write(word + " " + message.content[0].text + "\n")
        return None


def generate_image_from_word_info(word_info: Dict[str, str], generator: str = "flux.1-dev") -> BytesIO:
    """
    Generate an image from a given word's definition and example sentence.

    Args:
        word_info (Dict[str, str]): A dictionary containing the word's definition and example sentence under the keys 'definition' and 'example'.
        generator (str): The generator to use to generate the image. It can be "flux.1-dev" or "dall-e-3".
    Returns:
        BytesIO: A BytesIO object containing the generated image.
    """
    image_prompt = (
        f"Please make a picture of the word \"{word_info['word']}\", so a 5 year old can understand what the word means. "
        f"The definition of the word is: \"{word_info['definition']}\". "
        f"The example of the word is: \"{word_info['example']}\"."
    )

    if generator == "flux.1-dev":
        response = openai_client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (f"You are a prompt engineer. I want you to convert and expand a prompt "
                                     f"for use in text to image generation services which are based on Google T5 encoder and Flux model. "
                                     f"Convert following prompt to natural language, creating an expanded and detailed prompt "
                                     f"with detailed descriptions of subjects, scene and image quality while keeping the same key points. "
                                     f"The final output should combine all these elements into a cohesive, detailed prompt "
                                     f"that accurately reflects the image and should be converted into single paragraph "
                                     f"to give the best possible result. "
                                     f"\n\nThe prompt is: \"{image_prompt}\"")
                        }
                    ]
                }
            ],
            response_format={
                "type": "text"
            },
            temperature=0,
            max_completion_tokens=512,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=0
        )
        expanded_prompt = response.choices[0].message.content
        print("Expanded image prompt: " + expanded_prompt)
        image_data = generate_image(expanded_prompt)
    else:
        response = openai_client.images.generate(
            model="dall-e-3",
            prompt=image_prompt,
            size="1024x1024",
            quality="standard",
            n=1,
            # response_format="b64_json",
    )
        image_url_response = requests.get(response.data[0].url)
        image_data = BytesIO(image_url_response.content)

    return image_data

def generate_slide(word_info: Dict[str, str], image_data: BytesIO) -> str:
    """
    Generate a slide from a given word info and image data.

    Args:
        word_info (Dict[str, str]): A dictionary containing the word's definition and example sentence under the keys 'definition' and 'example'.
        image_data (BytesIO): A BytesIO object containing the generated image.

    Returns:
        str: The path to the generated slide.
    """
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    # set title
    title = slide.shapes.title
    title.top = Inches(0.5)
    title.left = Inches(1)
    title.height = Inches(1)
    title.width = Inches(8)
    title.text = word_info["word"]

    # add image to slide
    pic = slide.shapes.add_picture(image_data, left=Inches(1), top=Inches(1.5), width=Inches(5), height=Inches(5))

    # Add word definition
    txBox = slide.shapes.add_textbox(left=Inches(6), top=Inches(1.5), width=Inches(4), height=Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.text = word_info["definition"]

    # Add word example
    txBox = slide.shapes.add_textbox(left=Inches(6), top=Inches(4), width=Inches(4), height=Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.font.size = Pt(22)
    p.text = word_info["example"]

    # save pptx file
    output_file_path = cache_dir + 'temp.pptx'
    prs.save(output_file_path)

    return output_file_path

def convert_slide_to_image(slide_path: str) -> str:
    """
    Convert a PowerPoint slide to an image.

    Args:
        slide_path (str): The path to the PowerPoint slide.

    Returns:
        str: The path to the generated image.
    """
    pdf_path = cache_dir + "temp.pdf"

    # Convert PPTX to PDF using LibreOffice
    # TODO: check the os and use the correct path
    # MAC: /Applications/LibreOffice.app/Contents/MacOS/soffice
    # LINUX: soffice
    subprocess.run(["soffice", "--headless", "--convert-to", "pdf", slide_path, "--outdir", cache_dir], check=True)

    # Convert PDF to PNG images using pdftoppm (part of poppler-utils)
    subprocess.run(["pdftoppm", "-png", "-singlefile", pdf_path, cache_dir + "temp"], check=True)

    os.remove(pdf_path)

    return cache_dir + "temp.png"


def combine_mp3s(file_paths: List[str], output_path: str, pause_duration: int = 2000) -> str:
    """Combines multiple MP3 files into one.

    Args:
        file_paths (List[str]): A list of file paths to the MP3 files to combine.
        output_path (str): The file path to save the combined MP3 file.
        pause_duration (int): The duration of the pause between the audio files.

    Returns:
        str: The path to the generated audio file.
    """
    pause = AudioSegment.silent(duration=pause_duration)
    combined_audio = AudioSegment.empty()
    for i, file_path in enumerate(file_paths):
        audio = AudioSegment.from_mp3(file_path)
        combined_audio += audio
        combined_audio += pause
    combined_audio.export(output_path, format="mp3")


def generate_audio_file(word_info: Dict[str, str]) -> str:
    """
    Generate a audio file for a given word info.

    Args:
        word_info (Dict[str, str]): A dictionary containing the word's definition and example sentence under the keys 'definition' and 'example'.

    Returns:
        str: The path to the generated audio file.
    """
    word_speech_file_path = cache_dir + "word_speech.mp3"  
    with openai_client.audio.speech.with_streaming_response.create(
        model="tts-1",
        voice="shimmer",
        input= "The word is \"" + word_info["word"] + "\" ...",
    ) as response:
        response.stream_to_file(word_speech_file_path)

    word_definition_speech_file_path = cache_dir + "word_definition_speech.mp3"    
    with openai_client.audio.speech.with_streaming_response.create(
        model="tts-1",
        voice="shimmer",
        input= word_info["definition"],
    ) as response:
        response.stream_to_file(word_definition_speech_file_path)

    word_example_speech_file_path = cache_dir + "word_example_speech.mp3"
    with openai_client.audio.speech.with_streaming_response.create(
        model="tts-1",
        voice="nova",
        input= word_info["example"],
    ) as response:
        response.stream_to_file(word_example_speech_file_path)

    file_paths = [word_speech_file_path, word_definition_speech_file_path, word_example_speech_file_path]
    output_path = cache_dir + "combined_audio.mp3"
    combine_mp3s(file_paths, output_path)

    os.remove(word_speech_file_path)
    os.remove(word_definition_speech_file_path)
    os.remove(word_example_speech_file_path)

    return output_path


def generate_video_file(image_path: str, audio_path: str, output_path: str) -> str:
    """
    Generate a video file from a given image and audio file.

    Args:
        image_path (str): The path to the image file.
        audio_path (str): The path to the audio file.
        output_path (str): The path to save the generated video file.        
    """
    audio_clip = AudioFileClip(audio_path)
    image_clip = ImageClip(image_path).set_duration(audio_clip.duration)
    
    video_clip = image_clip.set_audio(audio_clip)
    video_clip.write_videofile(output_path, fps=24)


def generate_video_for_word(word: str, output_path: str):
    logging.debug("Starting video generation process...")
    
    logging.debug("Generating word information...")
    word_info = generate_word_info(word)
    if word_info is None:
        logging.debug("Failed to generate word info for word: " + word)
        return False
    logging.debug(f"Generated word info for '{word_info['word']}'")
    
    logging.debug("Generating image from word info...")
    image_data = generate_image_from_word_info(word_info)
    logging.debug("Image generated successfully")
    
    logging.debug("Creating PowerPoint slide...")
    slide_path = generate_slide(word_info, image_data)
    logging.debug(f"Slide created at {slide_path}")
    
    logging.debug("Converting slide to image...")
    image_path = convert_slide_to_image(slide_path)
    logging.debug(f"Slide converted to image at {image_path}")
    
    logging.debug("Generating audio file...")
    audio_path = generate_audio_file(word_info)
    logging.debug(f"Audio generated at {audio_path}")
    
    logging.debug("Creating final video...")
    generate_video_file(image_path, audio_path, output_path)
    logging.debug(f"Video created successfully at {output_path}")
    
    logging.debug("Deleting temporary files...")
    os.remove(slide_path)
    os.remove(image_path)
    os.remove(audio_path)
    logging.debug("Temporary files deleted successfully")

    logging.debug("Process completed!")

    return True

def combine_videos(video_paths: List[str], output_path: str):
    """
    Combine multiple video files into one.

    Args:
        video_paths (List[str]): A list of file paths to the video files to combine.
        output_path (str): The path to save the combined video file.
    """
    video_clips = [VideoFileClip(path) for path in video_paths]
    final_clip = concatenate_videoclips(video_clips, method="compose")
    final_clip.write_videofile(output_path, codec="libx264")

def combine_videos_from_cache_files(video_ids: List[str], output_path: str):
    """
    Combine multiple video files from cache files into one.
    """
    video_paths = []
    for video_to_include in video_ids:
        for file_path in os.listdir(cache_dir):
            file_abs_path = cache_dir + file_path
            if os.path.isfile(file_abs_path) and file_path.startswith(video_to_include+'_'):
                video_paths.append(file_abs_path)
    combine_videos(video_paths, output_path)


if __name__ == "__main__":
    # count is used for counting and skipping
    count = 10
    videos_to_include = [] 
    video_paths = []
    for video_to_include in videos_to_include:
        for file_path in os.listdir(cache_dir):
            file_abs_path = cache_dir + file_path
            if os.path.isfile(file_abs_path) and file_path.startswith(video_to_include+'_'):
                video_paths.append(file_abs_path)

    for i, word in enumerate(tqdm(word_list[count:20])):
        video_path = cache_dir + f"{count+1:02d}_" + word + ".mp4"
        success = generate_video_for_word(word, video_path)
        count += 1
        if success:
            video_paths.append(video_path)
        if count % 10 == 0:
            combine_videos(video_paths, output_dir + f"combined_{count-9}-{count}.mp4")
            video_paths = []

    # combine_videos_from_cache_files(["11", "12", "13", "14", "15", "16", "17", "18", "19", "20"], output_dir + "combined_11-20.mp4")